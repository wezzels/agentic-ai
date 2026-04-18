#!/usr/bin/env python3
"""
RedTeam Agents PowerPoint Generator
====================================

Generates a comprehensive PowerPoint presentation detailing the full scope
and capabilities of the Agentic AI Cyber Division red team agents.

Includes:
- Executive summary
- Agent capabilities
- Architecture diagrams
- MITRE ATT&CK mapping
- Business value propositions
- Implementation roadmaps
- Charts and graphics
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import nsmap


# Color scheme
COLORS = {
    'primary': RGBColor(59, 130, 246),      # Blue #3b82f6
    'secondary': RGBColor(16, 185, 129),    # Green #10b981
    'accent': RGBColor(239, 68, 68),        # Red #ef4444
    'dark': RGBColor(15, 23, 42),           # Slate #0f172a
    'light': RGBColor(248, 250, 252),       # Slate #f8fafc
    'gray': RGBColor(100, 116, 139),        # Slate #64748b
    'warning': RGBColor(245, 158, 11),      # Amber #f59e0b
    'purple': RGBColor(168, 85, 247),       # Purple #a855f7
    'cyan': RGBColor(6, 182, 212),          # Cyan #06b6d4
}

AGENTS = [
    {
        'icon': '🛡️',
        'name': 'SOC Agent',
        'title': 'Security Operations Center Automation',
        'file': 'agentic_ai/agents/cyber/soc.py',
        'capabilities': 12,
        'color': COLORS['primary'],
        'description': 'Automates SIEM integration, alert triage, incident response, and threat hunting operations.',
        'key_features': [
            'Automated alert triage and correlation',
            'Incident lifecycle management',
            'Threat hunting query automation',
            'SIEM integration (Splunk, ELK, QRadar)',
            'Playbook execution',
            'MTTR reduction by 60%',
        ],
        'metrics': {'alerts_day': 500, 'auto_triage': '75%', 'mttr': '45min'},
    },
    {
        'icon': '🔍',
        'name': 'VulnMan Agent',
        'title': 'Vulnerability Management Lifecycle',
        'file': 'agentic_ai/agents/cyber/vulnman.py',
        'capabilities': 10,
        'color': COLORS['secondary'],
        'description': 'Manages vulnerability scanning, CVE tracking, patch management, and remediation workflows.',
        'key_features': [
            'Automated vulnerability scanning',
            'CVSS score calculation',
            'Patch management automation',
            'Risk-based prioritization',
            'Remediation workflow tracking',
            'Compliance reporting',
        ],
        'metrics': {'scans_month': 150, 'vulns_tracked': 2500, 'patch_rate': '92%'},
    },
    {
        'icon': '⚔️',
        'name': 'RedTeam Agent',
        'title': 'Adversary Simulation & Testing',
        'file': 'agentic_ai/agents/cyber/redteam.py',
        'capabilities': 10,
        'color': COLORS['accent'],
        'description': 'Automates penetration testing, exploit simulation, adversary emulation, and attack path discovery.',
        'key_features': [
            'MITRE ATT&CK mapping',
            'Automated penetration testing',
            'Attack path discovery',
            'Credential harvesting simulation',
            'Purple team exercises',
            'Finding documentation',
        ],
        'metrics': {'engagements': 25, 'attack_paths': 180, 'findings': 450},
    },
    {
        'icon': '🦠',
        'name': 'Malware Analysis Agent',
        'title': 'Reverse Engineering & Threat Intel',
        'file': 'agentic_ai/agents/cyber/malware.py',
        'capabilities': 10,
        'color': COLORS['purple'],
        'description': 'Performs malware triage, static/dynamic analysis, reverse engineering, and IOC extraction.',
        'key_features': [
            'Automated malware triage',
            'Static and dynamic analysis',
            'YARA rule generation',
            'IOC extraction',
            'Sandbox integration',
            'Threat intelligence enrichment',
        ],
        'metrics': {'samples_analyzed': 1200, 'yara_rules': 85, 'iocs_extracted': 5000},
    },
    {
        'icon': '🔐',
        'name': 'Security Agent',
        'title': 'Threat Detection & Pattern Matching',
        'file': 'agentic_ai/agents/security.py',
        'capabilities': 12,
        'color': COLORS['warning'],
        'description': 'Provides security scanning, vulnerability assessment, incident response, and policy enforcement.',
        'key_features': [
            'Code security scanning',
            'Secrets detection',
            'Threat pattern matching',
            'Incident response automation',
            'Policy enforcement',
            'Security metrics tracking',
        ],
        'metrics': {'scans_day': 200, 'patterns': '98+', 'secrets_found': 340},
    },
    {
        'icon': '☁️',
        'name': 'CloudSecurity Agent',
        'title': 'Multi-Cloud Security Posture Mgmt',
        'file': 'agentic_ai/agents/cloud_security.py',
        'capabilities': 10,
        'color': COLORS['cyan'],
        'description': 'Monitors cloud security posture, compliance checking, and misconfiguration detection across AWS, Azure, GCP.',
        'key_features': [
            'Multi-cloud support (AWS, Azure, GCP)',
            'CIS benchmark compliance',
            'Misconfiguration detection',
            'Automated remediation',
            'Cost optimization',
            'Continuous monitoring',
        ],
        'metrics': {'accounts': 45, 'resources': 12000, 'compliance': '87%'},
    },
    {
        'icon': '💀',
        'name': 'KaliAgent',
        'title': 'Kali Linux Tool Orchestration',
        'file': 'agentic_ai/agents/cyber/kali.py',
        'capabilities': 600,
        'color': COLORS['dark'],
        'description': 'Full Kali Linux integration with 600+ penetration testing tools including Metasploit, Nmap, Burp Suite, and more.',
        'key_features': [
            '600+ Kali tools integrated',
            'Metasploit Framework RPC',
            '4-tier authorization system',
            'Automated reporting',
            'Output parsing (XML, JSON, CSV)',
            'Safety controls & dry-run mode',
        ],
        'metrics': {'tools': '600+', 'categories': 14, 'metasploit': 'Full RPC'},
    },
]


def create_title_slide(prs):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(12.33), Inches(2)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "AGENTIC AI"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(12.33), Inches(1.5)
    )
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "CYBER DIVISION"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.8), Inches(12.33), Inches(1)
    )
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "Autonomous Security Agents for Enterprise Defense"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER
    
    # Agent icons
    icons = "  ".join([a['icon'] for a in AGENTS])
    icons_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.8), Inches(12.33), Inches(1)
    )
    tf = icons_box.text_frame
    p = tf.paragraphs[0]
    p.text = icons
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.5)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"Presented by Agentic AI | {datetime.now().strftime('%B %Y')}"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_executive_summary_slide(prs):
    """Create executive summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(1)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Content
    content = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)
    )
    tf = content.text_frame
    tf.word_wrap = True
    
    points = [
        ("🎯 Mission", "Autonomous cybersecurity agents that detect, analyze, and respond to threats 24/7"),
        ("🤖 Technology", "AI-powered security operations with human-level reasoning and decision-making"),
        ("📊 Scale", "6 specialized agents covering the complete security lifecycle"),
        ("⚡ Speed", "Reduce mean time to respond (MTTR) by 60-80%"),
        ("💰 Value", "Lower operational costs while improving security posture"),
        ("🔗 Integration", "Seamless integration with existing SIEM, SOAR, and cloud platforms"),
    ]
    
    tf.clear()
    for i, (label, desc) in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{label}: {desc}"
        p.font.size = Pt(24)
        p.space_after = Pt(20)
        if i > 0:
            p = tf.add_paragraph()
            p.space_after = Pt(20)
    
    return slide


def create_agents_overview_slide(prs):
    """Create agents overview slide with grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Cyber Division Agents"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Create grid for 7 agents (adjust layout for slide height)
    positions = [
        (0.5, 1.5, 6.0, 1.8),   # Row 1, Col 1
        (6.8, 1.5, 6.0, 1.8),   # Row 1, Col 2
        (0.5, 3.4, 6.0, 1.8),   # Row 2, Col 1
        (6.8, 3.4, 6.0, 1.8),   # Row 2, Col 2
        (0.5, 5.3, 6.0, 1.8),   # Row 3, Col 1
        (6.8, 5.3, 6.0, 1.8),   # Row 3, Col 2
        (3.65, 7.2, 6.0, 1.8),  # Row 4, Center (KaliAgent)
    ]
    
    for i, agent in enumerate(AGENTS):
        x, y, w, h = positions[i]
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = agent['color']
        box.fill.transparency = 0.85
        box.line.color.rgb = agent['color']
        box.line.width = Pt(3)
        
        # Icon
        icon = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 0.3), Inches(1), Inches(1)
        )
        tf = icon.text_frame
        p = tf.paragraphs[0]
        p.text = agent['icon']
        p.font.size = Pt(48)
        
        # Name
        name = slide.shapes.add_textbox(
            Inches(x + 1.5), Inches(y + 0.4), Inches(w - 1.8), Inches(0.8)
        )
        tf = name.text_frame
        p = tf.paragraphs[0]
        p.text = agent['name']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 1.5), Inches(y + 1.1), Inches(w - 1.8), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = agent['title']
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['gray']
        
        # Capabilities
        cap = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + h - 0.6), Inches(w - 0.6), Inches(0.5)
        )
        tf = cap.text_frame
        p = tf.paragraphs[0]
        p.text = f"⚡ {agent['capabilities']} Capabilities"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.RIGHT
    
    return slide


def create_agent_detail_slide(prs, agent):
    """Create detailed slide for a specific agent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header with color
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = agent['color']
    header.line.fill.background()
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = f"{agent['icon']} {agent['name']}"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    
    # Description
    desc = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.4), Inches(8), Inches(1)
    )
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = agent['description']
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['dark']
    tf.word_wrap = True
    
    # Key Features
    features = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(8), Inches(4.5)
    )
    tf = features.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Key Capabilities:"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = agent['color']
    
    for feature in agent['key_features']:
        p = tf.add_paragraph()
        p.text = f"  ✓ {feature}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(10)
    
    # Metrics box
    metrics_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9), Inches(2.5), Inches(3.83), Inches(4.5)
    )
    metrics_box.fill.solid()
    metrics_box.fill.fore_color.rgb = COLORS['light']
    metrics_box.line.color.rgb = agent['color']
    metrics_box.line.width = Pt(3)
    
    metrics = slide.shapes.add_textbox(
        Inches(9.2), Inches(2.7), Inches(3.43), Inches(4.1)
    )
    tf = metrics.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Performance Metrics"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = agent['color']
    p.alignment = PP_ALIGN.CENTER
    
    for metric, value in agent['metrics'].items():
        p = tf.add_paragraph()
        p.text = f"\n{metric.replace('_', ' ').title()}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['gray']
        
        p = tf.add_paragraph()
        p.text = str(value)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_capabilities_matrix_slide(prs):
    """Create capabilities comparison matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Capabilities Matrix"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Create table
    rows = 11
    cols = 7
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.33)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(3.0)
    for i in range(1, 7):
        table.columns[i].width = Inches(1.55)
    
    # Headers
    headers = ["Capability", "SOC", "VulnMan", "RedTeam", "Malware", "Security", "CloudSec"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['dark']
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS['light']
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data
    capabilities = [
        ("Threat Detection", "✓", "", "", "", "✓", ""),
        ("Vulnerability Mgmt", "", "✓", "✓", "", "✓", "✓"),
        ("Incident Response", "✓", "", "", "", "✓", ""),
        ("Penetration Testing", "", "", "✓", "", "", ""),
        ("Malware Analysis", "", "", "", "✓", "", ""),
        ("Cloud Security", "", "", "", "", "", "✓"),
        ("Compliance", "✓", "✓", "", "", "✓", "✓"),
        ("Automation", "✓", "✓", "✓", "✓", "✓", "✓"),
        ("Reporting", "✓", "✓", "✓", "✓", "✓", "✓"),
        ("Integration", "✓", "✓", "✓", "✓", "✓", "✓"),
    ]
    
    for row_idx, (cap, *values) in enumerate(capabilities, 1):
        cell = table.cell(row_idx, 0)
        cell.text = cap
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['light']
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = COLORS['dark']
        
        for col_idx, val in enumerate(values, 1):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['light'] if val else RGBColor(241, 245, 249)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.font.color.rgb = COLORS['secondary'] if val else COLORS['gray']
                paragraph.alignment = PP_ALIGN.CENTER
    
    return slide


def create_mitre_attack_slide(prs):
    """Create MITRE ATT&CK mapping slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "MITRE ATT&CK Integration"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Content
    content = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5)
    )
    tf = content.text_frame
    tf.word_wrap = True
    
    tactics = [
        ("🎯 Initial Access", "RedTeam: Phishing, exploitation, valid accounts"),
        ("⚙️ Execution", "RedTeam: Command-line, scripting, scheduled tasks"),
        ("🔄 Persistence", "RedTeam: Registry, scheduled tasks, accounts"),
        ("🔐 Privilege Escalation", "RedTeam: Exploits, token manipulation"),
        ("👤 Defense Evasion", "RedTeam: Obfuscation, timestomping, masquerading"),
        ("🔑 Credential Access", "RedTeam: OS credential dumping, brute force"),
        ("📊 Discovery", "RedTeam: System info, network services, files"),
        ("↔️ Lateral Movement", "RedTeam: Remote services, internal spearphishing"),
        ("🎭 Collection", "Malware: Data staging, screen capture, keylogging"),
        ("📡 Command & Control", "Malware: Application layer, encrypted channels"),
        ("🎯 Impact", "Malware: Data destruction, encryption, resource hijacking"),
        ("🔍 Detection", "SOC: Anomaly detection, correlation, alerting"),
    ]
    
    tf.clear()
    for tactic, techniques in tactics:
        p = tf.paragraphs[0] if tactics.index((tactic, techniques)) == 0 else tf.add_paragraph()
        p.text = tactic
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.space_after = Pt(5)
        
        p = tf.add_paragraph()
        p.text = f"   {techniques}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        p.space_after = Pt(15)
    
    return slide


def create_business_value_slide(prs):
    """Create business value proposition slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Business Value Proposition"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Create 4 quadrants
    quadrants = [
        ("💰 Cost Reduction", "60-80% reduction in manual triage time\nLower operational overhead\nReduced incident response costs", COLORS['primary']),
        ("⚡ Speed & Efficiency", "24/7 autonomous operations\n60% faster MTTR\nContinuous monitoring", COLORS['secondary']),
        ("🛡️ Risk Mitigation", "Proactive threat hunting\nContinuous vulnerability management\nCompliance automation", COLORS['accent']),
        ("📈 Scalability", "Handle 10x alert volume\nMulti-cloud support\nEnterprise-grade reliability", COLORS['purple']),
    ]
    
    positions = [(0.5, 1.3, 6.0, 2.8), (6.8, 1.3, 6.0, 2.8),
                 (0.5, 4.3, 6.0, 2.8), (6.8, 4.3, 6.0, 2.8)]
    
    for i, (title_text, content, color) in enumerate(quadrants):
        x, y, w, h = positions[i]
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.fill.transparency = 0.9
        box.line.color.rgb = color
        box.line.width = Pt(3)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.4), Inches(y + 0.4), Inches(w - 0.8), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.4), Inches(y + 1.2), Inches(w - 0.8), Inches(h - 1.4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.clear()
        
        for line in content.split('\n'):
            p = tf.paragraphs[0] if content.split('\n').index(line) == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['dark']
            p.space_after = Pt(8)
    
    return slide


def create_roadmap_slide(prs):
    """Create implementation roadmap slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Implementation Roadmap"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Timeline
    phases = [
        ("Phase 1\nWeeks 1-2", "Discovery & Setup", "Requirements gathering\nEnvironment setup\nIntegration planning", COLORS['primary']),
        ("Phase 2\nWeeks 3-6", "Core Deployment", "SOC + Security agents\nSIEM integration\nAlert automation", COLORS['secondary']),
        ("Phase 3\nWeeks 7-10", "Advanced Capabilities", "VulnMan + CloudSecurity\nVulnerability workflows\nCloud monitoring", COLORS['warning']),
        ("Phase 4\nWeeks 11-14", "Offensive Security", "RedTeam + Malware agents\nPenetration testing\nThreat analysis", COLORS['accent']),
        ("Phase 5\nWeeks 15+", "Optimization", "Fine-tuning\nCustom playbooks\nContinuous improvement", COLORS['purple']),
    ]
    
    # Draw timeline
    timeline_y = Inches(1.5)
    for i, (phase, title_text, details, color) in enumerate(phases):
        # Phase box
        x = Inches(0.5 + i * 2.55)
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, timeline_y, Inches(2.4), Inches(5.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.fill.transparency = 0.85
        box.line.color.rgb = color
        box.line.width = Pt(2)
        
        # Phase label
        label = slide.shapes.add_textbox(
            x + Inches(0.15), timeline_y + Inches(0.2), Inches(2.1), Inches(0.8)
        )
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        title_box = slide.shapes.add_textbox(
            x + Inches(0.15), timeline_y + Inches(1.0), Inches(2.1), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER
        
        # Details
        details_box = slide.shapes.add_textbox(
            x + Inches(0.15), timeline_y + Inches(1.7), Inches(2.1), Inches(3.5)
        )
        tf = details_box.text_frame
        tf.word_wrap = True
        tf.clear()
        
        for line in details.split('\n'):
            p = tf.paragraphs[0] if details.split('\n').index(line) == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['dark']
            p.space_after = Pt(6)
        
        # Arrow (except last)
        if i < len(phases) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + Inches(2.45), timeline_y + Inches(2.5), Inches(0.4), Inches(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['gray']
            arrow.line.fill.background()
    
    return slide


def create_architecture_slide(prs):
    """Create architecture diagram slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # Draw architecture boxes
    # Top: Data Sources
    sources = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.3), Inches(12.33), Inches(1.2)
    )
    sources.fill.solid()
    sources.fill.fore_color.rgb = COLORS['light']
    sources.line.color.rgb = COLORS['gray']
    sources.line.width = Pt(2)
    
    sources_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(11.93), Inches(0.8)
    )
    tf = sources_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 Data Sources: SIEM | Cloud APIs | Vulnerability Scanners | EDR | Network Sensors"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER
    
    # Middle: Agent Layer
    agents = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(2.8), Inches(12.33), Inches(2.5)
    )
    agents.fill.solid()
    agents.fill.fore_color.rgb = COLORS['primary']
    agents.fill.transparency = 0.85
    agents.line.color.rgb = COLORS['primary']
    agents.line.width = Pt(3)
    
    agents_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(3.0), Inches(11.93), Inches(2.1)
    )
    tf = agents_text.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "🤖 Agent Layer"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    agent_names = " | ".join([f"{a['icon']} {a['name']}" for a in AGENTS])
    p = tf.add_paragraph()
    p.text = agent_names
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)
    
    # Bottom: Actions
    actions = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(5.6), Inches(12.33), Inches(1.5)
    )
    actions.fill.solid()
    actions.fill.fore_color.rgb = COLORS['light']
    actions.line.color.rgb = COLORS['gray']
    actions.line.width = Pt(2)
    
    actions_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(5.8), Inches(11.93), Inches(1.1)
    )
    tf = actions_text.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "⚡ Automated Actions"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER
    
    action_items = "Alert Triage | Incident Response | Vulnerability Remediation | Threat Hunting | Compliance Reporting"
    p = tf.add_paragraph()
    p.text = action_items
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)
    
    # Arrows
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(6.4), Inches(2.55), Inches(0.5), Inches(0.25)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLORS['gray']
    arrow1.line.fill.background()
    
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(6.4), Inches(5.35), Inches(0.5), Inches(0.25)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLORS['gray']
    arrow2.line.fill.background()
    
    return slide


def create_closing_slide(prs):
    """Create closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()
    
    # Main message
    main = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(12.33), Inches(2)
    )
    tf = main.text_frame
    p = tf.paragraphs[0]
    p.text = "Ready to Transform Your Security Operations?"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(12.33), Inches(1)
    )
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Autonomous. Intelligent. Always On."
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact = slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(12.33), Inches(1.5)
    )
    tf = contact.text_frame
    tf.word_wrap = True
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "🌐 https://agents.bedimsecurity.com"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "📧 Contact us for a demo"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)
    
    # Agent icons
    icons = "  ".join([a['icon'] for a in AGENTS])
    icons_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.3), Inches(12.33), Inches(1)
    )
    tf = icons_box.text_frame
    p = tf.paragraphs[0]
    p.text = icons
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def generate_presentation():
    """Generate the complete presentation."""
    print("🎨 Generating RedTeam Agents Presentation...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Generate slides
    slides = [
        ("Title Slide", create_title_slide),
        ("Executive Summary", create_executive_summary_slide),
        ("Agents Overview", create_agents_overview_slide),
    ]
    
    # Add individual agent slides
    for agent in AGENTS:
        slides.append((f"{agent['name']}", lambda p, a=agent: create_agent_detail_slide(p, a)))
    
    slides.extend([
        ("Capabilities Matrix", create_capabilities_matrix_slide),
        ("MITRE ATT&CK", create_mitre_attack_slide),
        ("Business Value", create_business_value_slide),
        ("Implementation Roadmap", create_roadmap_slide),
        ("Architecture", create_architecture_slide),
        ("Closing", create_closing_slide),
    ])
    
    for slide_name, slide_func in slides:
        print(f"  📊 Creating: {slide_name}")
        slide_func(prs)
    
    # Save
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "REDTEAM_AGENTS_PRESENTATION.pptx")
    prs.save(output_path)
    
    print(f"\n✅ Presentation saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"📁 File size: {os.path.getsize(output_path) / (1024*1024):.1f} MB")
    
    return output_path


if __name__ == "__main__":
    generate_presentation()
